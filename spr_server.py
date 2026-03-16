#!/usr/bin/env python3
"""
spr_server.py - SPR Analysis Server (Upload-based)

Starts a web server with a file upload interface.
Users upload SPR data files, the server auto-detects MCK/SCK format,
runs the C fitting engine, and displays an interactive dashboard.

Usage:
    python3 spr_server.py [--port 8765]
"""

import subprocess
import json
import sys
import os
import io
import tempfile
import struct
import base64
from http.server import HTTPServer, BaseHTTPRequestHandler
from urllib.parse import urlparse

# ── Configuration ──────────────────────────────────────────────────

C_BINARY = "./spr_fit"
C_SOURCES = ["spr_fit_main.c", "spr_io.c", "spr_models.c", "spr_optim.c"]
C_HEADERS = ["spr_types.h", "spr_io.h", "spr_models.h", "spr_optim.h"]

CYCLE_COLORS = [
    "#1f77b4", "#ff7f0e", "#2ca02c", "#d62728",
    "#9467bd", "#8c564b", "#e377c2", "#7f7f7f",
    "#bcbd22", "#17becf"
]

UPLOAD_DIR = None  # Set in main()


def compile_fitter():
    """Compile the C fitting engine if needed."""
    all_files = C_SOURCES + C_HEADERS
    existing = [f for f in all_files if os.path.exists(f)]
    if not existing:
        print("No C source files found!")
        return False
    src_mtime = max(os.path.getmtime(f) for f in existing)
    if os.path.exists(C_BINARY) and os.path.getmtime(C_BINARY) >= src_mtime:
        return True
    print("Compiling spr_fit ...")
    ret = subprocess.run(["make"], capture_output=True, text=True)
    if ret.returncode != 0:
        print(f"Compilation failed:\n{ret.stderr}")
        return False
    print("Compilation successful.")
    return True


def run_fitter(filepath, ref_path=None, model="langmuir", exclude=None,
               assoc_end=None, ri=False, drift=False, tc=False):
    """Run spr_fit on a data file and return parsed JSON result."""
    cmd = [C_BINARY, filepath]
    if ref_path and os.path.exists(ref_path):
        cmd += ["--ref", ref_path]
    if model != "langmuir":
        cmd += ["--model", model]
    if exclude:
        cmd += ["--exclude", ",".join(str(i) for i in exclude)]
    if assoc_end is not None:
        cmd += ["--assoc-end", str(assoc_end)]
    if ri:
        cmd.append("--ri")
    if drift:
        cmd.append("--drift")
    if tc:
        cmd.append("--tc")

    print(f"  Fitting {os.path.basename(filepath)} ...")
    ret = subprocess.run(cmd, capture_output=True, text=True, timeout=120)

    # Print fitting log (stderr)
    for line in ret.stderr.strip().split('\n'):
        if line:
            print(f"    {line}")

    if ret.returncode != 0:
        return None

    try:
        return json.loads(ret.stdout)
    except json.JSONDecodeError as e:
        print(f"    JSON parse error: {e}")
        return None


def parse_multipart(content_type, body):
    """Parse multipart/form-data body. Returns dict of field_name -> value.
    File fields get {'filename': ..., 'data': bytes}.
    Text fields get string value."""
    boundary = None
    for part in content_type.split(';'):
        part = part.strip()
        if part.startswith('boundary='):
            boundary = part[9:].strip('"')
            break
    if not boundary:
        return {}

    boundary_bytes = ('--' + boundary).encode()
    end_boundary = boundary_bytes + b'--'
    parts = body.split(boundary_bytes)

    result = {}
    for part in parts:
        if not part or part.strip() == b'' or part.strip() == b'--':
            continue
        if part.startswith(b'\r\n'):
            part = part[2:]
        if part.endswith(b'\r\n'):
            part = part[:-2]

        # Split headers from body
        header_end = part.find(b'\r\n\r\n')
        if header_end < 0:
            continue
        headers_raw = part[:header_end].decode('utf-8', errors='replace')
        data = part[header_end + 4:]
        if data.endswith(b'\r\n'):
            data = data[:-2]

        # Parse Content-Disposition
        name = None
        filename = None
        for line in headers_raw.split('\r\n'):
            if 'Content-Disposition' in line:
                for attr in line.split(';'):
                    attr = attr.strip()
                    if attr.startswith('name='):
                        name = attr[5:].strip('"')
                    elif attr.startswith('filename='):
                        filename = attr[9:].strip('"')

        if name:
            entry = {'filename': filename, 'data': data} if filename else data.decode('utf-8', errors='replace')
            if name in result:
                # Multiple values for same field name — convert to list
                if not isinstance(result[name], list):
                    result[name] = [result[name]]
                result[name].append(entry)
            else:
                result[name] = entry

    return result


def generate_pptx(datasets, dataset_idx=None):
    """Generate PPTX with 4 slides per dataset:
    1. Sensorgram, 2. Steady-State, 3. Residuals, 4. Parameters table."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import math

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    indices = range(len(datasets)) if dataset_idx is None else [dataset_idx]

    def add_slide_title(slide, title, subtitle=None):
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.15),
                                          Inches(12), Inches(0.7))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x2c, 0x3e, 0x50)
        if subtitle:
            sub = tf.add_paragraph()
            sub.text = subtitle
            sub.font.size = Pt(12)
            sub.font.color.rgb = RGBColor(0x7f, 0x8c, 0x8d)

    def fig_to_image(fig):
        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=200)
        plt.close(fig)
        buf.seek(0)
        return buf

    for idx in indices:
        d = datasets[idx]
        label = d.get('analyte', f'Dataset {idx+1}')
        model = d.get('model', 'langmuir')
        model_label = {'langmuir': '1:1 Langmuir',
                       'heterogeneous': 'Heterogeneous Ligand',
                       'twostate': 'Two-State'}.get(model, model)
        cycles = d.get('cycles', [])
        colors = CYCLE_COLORS
        # Build concentration → color mapping so replicates share colors
        unique_concs = []
        for cy in cycles:
            cn = round(cy.get('conc_nM', 0), 4)
            if cn not in unique_concs:
                unique_concs.append(cn)
        conc_color = {cn: colors[j % len(colors)] for j, cn in enumerate(unique_concs)}
        t_assoc = d.get('t_assoc_end', 60)
        ka, kd = d.get('ka', 0), d.get('kd', 0)
        KD = kd / ka if ka > 0 else 0
        Rmax = d.get('Rmax', 0)
        ss_kd = d.get('ss_KD_nM', 0)

        # ═══ Slide 1: Sensorgram ═══
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_title(slide, f'{label} — Sensorgram', f'Model: {model_label}')

        fig, ax = plt.subplots(figsize=(11, 5.5))
        for i, cy in enumerate(cycles):
            c = conc_color.get(round(cy.get('conc_nM', 0), 4), colors[0])
            time = cy.get('time', [])
            resp = cy.get('response', [])
            fitted = cy.get('fitted', [])
            skip = cy.get('skip', [])
            lbl = f"{cy.get('conc_nM', 0) * 1e-9:.2e} M"
            resp_plot = [r if not (skip and j < len(skip) and skip[j]) else float('nan')
                         for j, r in enumerate(resp)]
            ax.plot(time, resp_plot, color=c, linewidth=1.2, label=lbl)
            ax.plot(time, fitted, color=c, linewidth=1.8, linestyle='--', alpha=0.8)
        ax.axvline(t_assoc, color='#bdc3c7', linestyle=':', linewidth=1)
        ax.set_xlabel('Time (s)', fontsize=12)
        ax.set_ylabel('Response (RU)', fontsize=12)
        ax.set_title(f'{label} — {model_label}', fontsize=14)
        ax.legend(fontsize=9, loc='upper right', ncol=2)
        ax.grid(True, alpha=0.3)
        fig.tight_layout()
        slide.shapes.add_picture(fig_to_image(fig), Inches(0.5), Inches(1.0),
                                  Inches(12.3), Inches(6.2))

        # ═══ Slide 2: Steady-State Affinity ═══
        concs = d.get('concentrations', [])
        ss_req = d.get('ss_req', [])
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_title(slide, f'{label} — Steady-State Affinity',
                        f'KD = {ss_kd * 1e-9:.3e} M' if ss_kd > 0 else 'N/A')

        if concs and ss_req and ss_kd > 0:
            fig, ax = plt.subplots(figsize=(10, 5.5))
            concs_M = [c * 1e-9 for c in concs]
            ax.scatter(concs_M, ss_req, color='#e74c3c', s=80, zorder=5,
                       label='Req (data)')
            ss_rmax = d.get('ss_Rmax', 1)
            min_c = min(concs_M) * 0.3
            max_c = max(concs_M) * 2
            log_min = math.log10(min_c)
            log_max = math.log10(max_c)
            curve_c = [10 ** (log_min + (log_max - log_min) * i / 200)
                       for i in range(201)]
            curve_r = [ss_rmax * (c * 1e9) / (ss_kd + c * 1e9) for c in curve_c]
            ax.plot(curve_c, curve_r, color='#2c3e50', linewidth=2,
                    label='Steady-state fit')
            ax.set_xscale('log')
            ax.set_xlabel('Concentration (M)', fontsize=12)
            ax.set_ylabel('Req (RU)', fontsize=12)
            ax.set_title(f'Steady-State Affinity: KD = {ss_kd * 1e-9:.3e} M, '
                         f'R\u00b2 = {d.get("ss_R2", 0):.4f}', fontsize=14)
            ax.legend(fontsize=10, loc='lower right')
            ax.grid(True, alpha=0.3)
            fig.tight_layout()
            slide.shapes.add_picture(fig_to_image(fig), Inches(0.5), Inches(1.0),
                                      Inches(12.3), Inches(6.2))
        else:
            txBox = slide.shapes.add_textbox(Inches(3), Inches(3),
                                              Inches(7), Inches(1))
            txBox.text_frame.paragraphs[0].text = 'Steady-state analysis not available for this dataset.'
            txBox.text_frame.paragraphs[0].font.size = Pt(16)
            txBox.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x95, 0xa5, 0xa6)

        # ═══ Slide 3: Residuals ═══
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_title(slide, f'{label} — Residuals',
                        f'RMS = {d.get("rms", 0):.4f} RU')

        fig, ax = plt.subplots(figsize=(11, 5.5))
        for i, cy in enumerate(cycles):
            c = conc_color.get(round(cy.get('conc_nM', 0), 4), colors[0])
            time = cy.get('time', [])
            resp = cy.get('response', [])
            fitted = cy.get('fitted', [])
            skip = cy.get('skip', [])
            lbl = f"{cy.get('conc_nM', 0) * 1e-9:.2e} M"
            res = [r - f if not (skip and j < len(skip) and skip[j]) else float('nan')
                   for j, (r, f) in enumerate(zip(resp, fitted))]
            ax.plot(time, res, color=c, linewidth=0.8, label=lbl)
        ax.axhline(0, color='#e74c3c', linewidth=1)
        ax.axvline(t_assoc, color='#bdc3c7', linestyle=':', linewidth=1)
        ax.set_xlabel('Time (s)', fontsize=12)
        ax.set_ylabel('Residual (RU)', fontsize=12)
        ax.set_title('Residuals (Data \u2212 Fit)', fontsize=14)
        ax.legend(fontsize=9, loc='upper right', ncol=2)
        ax.grid(True, alpha=0.3)
        fig.tight_layout()
        slide.shapes.add_picture(fig_to_image(fig), Inches(0.5), Inches(1.0),
                                  Inches(12.3), Inches(6.2))

        # ═══ Slide 4: Parameters Table ═══
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_slide_title(slide, f'{label} — Fitted Parameters',
                        f'Model: {model_label}')

        rows_data = []
        if model == 'langmuir':
            rows_data = [
                ('ka', f'{ka:.4e}', 'M\u207b\u00b9s\u207b\u00b9'),
                ('kd', f'{kd:.4e}', 's\u207b\u00b9'),
                ('KD', f'{KD:.4e}', 'M'),
                ('Rmax', f'{Rmax:.3f}', 'RU'),
            ]
        elif model == 'heterogeneous':
            ka2, kd2 = d.get('ka2', 0), d.get('kd2', 0)
            KD2 = kd2 / ka2 if ka2 > 0 else 0
            Rmax2 = d.get('Rmax2', 0)
            rows_data = [
                ('ka1', f'{ka:.4e}', 'M\u207b\u00b9s\u207b\u00b9'),
                ('kd1', f'{kd:.4e}', 's\u207b\u00b9'),
                ('KD1', f'{KD:.4e}', 'M'),
                ('Rmax1', f'{Rmax:.3f}', 'RU'),
                ('ka2', f'{ka2:.4e}', 'M\u207b\u00b9s\u207b\u00b9'),
                ('kd2', f'{kd2:.4e}', 's\u207b\u00b9'),
                ('KD2', f'{KD2:.4e}', 'M'),
                ('Rmax2', f'{Rmax2:.3f}', 'RU'),
            ]
        else:
            ka2, kd2 = d.get('ka2', 0), d.get('kd2', 0)
            rows_data = [
                ('ka1', f'{ka:.4e}', 'M\u207b\u00b9s\u207b\u00b9'),
                ('kd1', f'{kd:.4e}', 's\u207b\u00b9'),
                ('KD1', f'{KD:.4e}', 'M'),
                ('ka2', f'{ka2:.4e}', 's\u207b\u00b9'),
                ('kd2', f'{kd2:.4e}', 's\u207b\u00b9'),
                ('Rmax', f'{Rmax:.3f}', 'RU'),
            ]

        rows_data.append(('R\u00b2', f'{d.get("R2", 0):.6f}', ''))
        rows_data.append(('RMS', f'{d.get("rms", 0):.4f}', 'RU'))
        rows_data.append(('Data points', f'{d.get("n_points", 0)}', ''))

        if ss_kd > 0:
            rows_data.append(('', '', ''))  # spacer
            rows_data.append(('KD (SS)', f'{ss_kd * 1e-9:.4e}', 'M'))
            rows_data.append(('Rmax (SS)', f'{d.get("ss_Rmax", 0):.3f}', 'RU'))
            rows_data.append(('R\u00b2 (SS)', f'{d.get("ss_R2", 0):.6f}', ''))

        n_rows = len(rows_data) + 1
        tbl = slide.shapes.add_table(n_rows, 3, Inches(3), Inches(1.2),
                                      Inches(7), Inches(0.4 * n_rows)).table
        tbl.columns[0].width = Inches(2.5)
        tbl.columns[1].width = Inches(3.0)
        tbl.columns[2].width = Inches(1.5)

        for ci, txt in enumerate(['Parameter', 'Value', 'Unit']):
            cell = tbl.cell(0, ci)
            cell.text = txt
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.bold = True
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x34, 0x49, 0x5e)

        for ri, (param, val, unit) in enumerate(rows_data):
            for ci, txt in enumerate([param, val, unit]):
                cell = tbl.cell(ri + 1, ci)
                cell.text = txt
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(11)
                    if ci == 1:
                        p.font.name = 'Consolas'
            if param.startswith('KD'):
                for ci in range(3):
                    tbl.cell(ri + 1, ci).fill.solid()
                    tbl.cell(ri + 1, ci).fill.fore_color.rgb = RGBColor(0xEA, 0xFA, 0xF1)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()


# ── HTML Template ──────────────────────────────────────────────────

HTML_TEMPLATE = r"""<!DOCTYPE html>
<html lang="en">
<head>
<meta charset="UTF-8">
<title>SPR Kinetics Analysis</title>
<script src="https://cdn.plot.ly/plotly-latest.min.js"></script>
<style>
* { margin: 0; padding: 0; box-sizing: border-box; }
body { font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif;
       background: #f5f6fa; color: #2c3e50; padding: 20px; }
h1 { text-align: center; margin-bottom: 5px; font-size: 1.6em; }
.subtitle { text-align: center; color: #7f8c8d; margin-bottom: 20px; font-size: 0.9em; }

/* Upload panel */
.upload-panel { max-width: 700px; margin: 60px auto; background: white;
                border-radius: 12px; padding: 40px; box-shadow: 0 4px 20px rgba(0,0,0,0.1); }
.upload-panel h2 { margin-bottom: 20px; color: #2c3e50; }
.upload-group { margin-bottom: 20px; }
.upload-group label { display: block; font-weight: 600; margin-bottom: 6px; }
.upload-group input[type=file] { font-size: 0.95em; }
.upload-group select, .upload-group input[type=number] {
    padding: 6px 12px; border-radius: 6px; border: 1px solid #bdc3c7;
    font-size: 0.95em; background: white; }
.upload-group .help { color: #95a5a6; font-size: 0.82em; margin-top: 4px; }
.inline-group { display: flex; gap: 20px; flex-wrap: wrap; }
.inline-group > div { flex: 1; min-width: 150px; }
.upload-btn { padding: 12px 40px; border-radius: 8px; border: none;
              background: #3498db; color: white; font-size: 1.05em;
              font-weight: 600; cursor: pointer; transition: all 0.2s; }
.upload-btn:hover { background: #2980b9; }
.upload-btn:disabled { background: #bdc3c7; cursor: wait; }
.upload-status { margin-top: 15px; font-size: 0.9em; color: #7f8c8d;
                  min-height: 20px; }
.drop-zone { border: 2px dashed #bdc3c7; border-radius: 10px; padding: 25px;
              text-align: center; color: #95a5a6; transition: all 0.2s;
              margin-bottom: 10px; cursor: pointer; }
.drop-zone.dragover { border-color: #3498db; background: #eaf2fd; color: #3498db; }
.drop-zone input { display: none; }

/* Dashboard (hidden initially) */
#dashboard { display: none; }
.controls { display: flex; gap: 15px; align-items: center; justify-content: center;
             margin-bottom: 15px; flex-wrap: wrap; }
.controls label { font-weight: 600; }
.controls select { padding: 6px 12px; border-radius: 6px; border: 1px solid #bdc3c7;
                    font-size: 0.95em; background: white; }
.plot-container { background: white; border-radius: 10px; padding: 15px;
                   box-shadow: 0 2px 8px rgba(0,0,0,0.08); margin-bottom: 20px; }
.params-box { background: white; border-radius: 10px; padding: 20px;
               box-shadow: 0 2px 8px rgba(0,0,0,0.08); margin-bottom: 20px; }
.params-box h3 { margin-bottom: 12px; color: #2c3e50; }
table { border-collapse: collapse; width: 100%; font-size: 0.9em; }
th { background: #34495e; color: white; padding: 10px 12px; text-align: left; }
td { padding: 8px 12px; border-bottom: 1px solid #ecf0f1; }
tr:hover td { background: #f0f3f6; }
.highlight { background: #eafaf1 !important; font-weight: 600; }
.metric { font-family: 'SF Mono', 'Consolas', monospace; }
.two-col { display: grid; grid-template-columns: 1fr 1fr; gap: 20px; }
@media (max-width: 1000px) { .two-col { grid-template-columns: 1fr; } }
.toggle-row { display: flex; gap: 10px; justify-content: center; margin: 5px 0; }
.toggle-btn { padding: 5px 14px; border-radius: 5px; border: 1px solid #bdc3c7;
               background: white; cursor: pointer; font-size: 0.85em; }
.toggle-btn.active { background: #3498db; color: white; border-color: #3498db; }
.refit-btn { padding: 6px 20px; border-radius: 6px; border: 2px solid #c0392b;
             background: #e74c3c; color: white; cursor: pointer; font-size: 0.9em;
             font-weight: 600; transition: all 0.2s; }
.refit-btn:hover { background: #c0392b; }
.refit-btn:disabled { background: #bdc3c7; border-color: #95a5a6; cursor: wait; }
.export-btn { padding: 6px 20px; border-radius: 6px; border: 2px solid #27ae60;
              background: #2ecc71; color: white; cursor: pointer; font-size: 0.9em;
              font-weight: 600; transition: all 0.2s; }
.export-btn:hover { background: #27ae60; }
.export-btn:disabled { background: #bdc3c7; border-color: #95a5a6; cursor: wait; }
.new-upload-btn { padding: 6px 20px; border-radius: 6px; border: 2px solid #9b59b6;
                  background: #8e44ad; color: white; cursor: pointer; font-size: 0.9em;
                  font-weight: 600; transition: all 0.2s; }
.new-upload-btn:hover { background: #9b59b6; }
.cycle-toggle { display: flex; gap: 8px; align-items: center; flex-wrap: wrap;
                justify-content: center; margin: 8px 0; }
.cycle-cb { display: flex; align-items: center; gap: 3px; font-size: 0.85em;
            cursor: pointer; padding: 2px 6px; border-radius: 4px;
            border: 1px solid #ddd; transition: all 0.2s; }
.cycle-cb.excluded { opacity: 0.4; text-decoration: line-through; }
.cycle-cb input { cursor: pointer; }
.fit-log { font-family: 'SF Mono', 'Consolas', monospace; font-size: 0.8em;
           background: #2c3e50; color: #2ecc71; padding: 10px; border-radius: 6px;
           max-height: 150px; overflow-y: auto; margin-top: 10px; white-space: pre-wrap; }
</style>
</head>
<body>

<h1>SPR Kinetics Analysis</h1>
<p class="subtitle">Global Fit &mdash; Multi-Cycle / Single-Cycle Kinetics</p>

<!-- ═══ Upload Panel ═══ -->
<div id="uploadPanel" class="upload-panel">
  <h2>Upload SPR Data</h2>

  <div class="upload-group">
    <label>Data File(s)</label>
    <div class="drop-zone" id="dropZone" onclick="document.getElementById('dataFiles').click()">
      <p>Drag & drop files here, or click to browse</p>
      <p style="font-size:0.8em;margin-top:5px;">Supports Biacore MCK and SCK export formats (.txt)</p>
      <input type="file" id="dataFiles" multiple accept=".txt,.csv,.tsv">
    </div>
    <div id="fileList" style="font-size:0.85em;color:#2c3e50;"></div>
  </div>

  <div class="upload-group">
    <label>Reference File (optional, for double referencing)</label>
    <div class="drop-zone" id="refDropZone" onclick="document.getElementById('refFile').click()"
         style="padding:15px;">
      <p>Drag & drop reference file here, or click to browse</p>
      <p style="font-size:0.8em;margin-top:4px;">e.g. blank surface / control channel for subtraction</p>
      <input type="file" id="refFile" accept=".txt,.csv,.tsv">
    </div>
    <div id="refFileList" style="font-size:0.85em;color:#2c3e50;"></div>
  </div>

  <div class="upload-group">
    <label>Dataset Label</label>
    <input type="text" id="dataLabel" placeholder="e.g. Compound-A, Ch2"
           style="padding:6px 12px;border-radius:6px;border:1px solid #bdc3c7;width:100%;font-size:0.95em;">
    <p class="help">Optional label for this dataset</p>
  </div>

  <div class="inline-group">
    <div class="upload-group">
      <label>Binding Model</label>
      <select id="uploadModel">
        <option value="langmuir">1:1 Langmuir</option>
        <option value="heterogeneous">Heterogeneous Ligand</option>
        <option value="twostate">Two-State (Conformational Change)</option>
      </select>
    </div>
  </div>

  <div style="text-align:center;margin-top:20px;">
    <button class="upload-btn" id="uploadBtn" onclick="doUpload()">Analyze</button>
  </div>
  <div class="upload-status" id="uploadStatus"></div>
</div>

<!-- ═══ Dashboard ═══ -->
<div id="dashboard">

<div class="controls">
  <label>Dataset:</label>
  <select id="datasetSelect" onchange="onDatasetChange()"></select>
  <label>Model:</label>
  <select id="modelSelect">
    <option value="langmuir">1:1 Langmuir</option>
    <option value="heterogeneous">Heterogeneous Ligand</option>
    <option value="twostate">Two-State (Conformational Change)</option>
  </select>
  <details style="display:inline-block;margin-left:12px;vertical-align:middle;">
    <summary style="cursor:pointer;font-size:13px;color:#2c3e50;">Advanced</summary>
    <div style="margin-top:6px;font-size:13px;">
      <label style="margin-right:10px;"><input type="checkbox" id="enableRI"> Bulk RI</label>
      <label style="margin-right:10px;"><input type="checkbox" id="enableDrift"> Drift</label>
      <label><input type="checkbox" id="enableTC"> Mass Transport</label>
    </div>
  </details>
</div>

<div class="controls">
  <div class="toggle-row">
    <button class="toggle-btn active" id="btnData" onclick="toggleTrace('data')">Show Data</button>
    <button class="toggle-btn active" id="btnFit" onclick="toggleTrace('fit')">Show Fitted Curves</button>
  </div>
  <button class="refit-btn" id="refitBtn" onclick="doRefit()">Re-fit</button>
  <button class="export-btn" id="exportBtn" onclick="doExport()">Export PPTX</button>
  <button class="new-upload-btn" onclick="showUpload()">+ Add Data</button>
</div>

<div class="cycle-toggle" id="cycleToggles"></div>

<div class="two-col">
  <div class="plot-container">
    <div id="sensorgram" style="height:450px;"></div>
  </div>
  <div class="plot-container">
    <div id="residuals" style="height:450px;"></div>
  </div>
</div>

<div class="two-col">
  <div class="plot-container">
    <div id="isotherm" style="height:400px;"></div>
  </div>
  <div class="params-box">
    <h3>Fitted Parameters</h3>
    <table id="paramsTable">
      <thead><tr><th>Parameter</th><th>Value</th><th>Unit</th></tr></thead>
      <tbody id="paramsBody"></tbody>
    </table>
    <div class="fit-log" id="fitLog" style="display:none;"></div>
  </div>
</div>

<div class="params-box">
  <h3>Summary &mdash; All Datasets</h3>
  <div style="overflow-x:auto;">
  <table id="summaryTable">
    <thead><tr>
      <th>Dataset</th>
      <th>Mode</th>
      <th>Ref Sub</th>
      <th>k<sub>a</sub> (M<sup>-1</sup>s<sup>-1</sup>)</th>
      <th>k<sub>d</sub> (s<sup>-1</sup>)</th>
      <th>K<sub>D</sub> (M)</th>
      <th>R<sub>max</sub> (RU)</th>
      <th>R&sup2;</th>
      <th>Model</th>
    </tr></thead>
    <tbody id="summaryBody"></tbody>
  </table>
  </div>
</div>

</div><!-- /dashboard -->

<script>
const COLORS = __COLORS_JSON__;
const SERVER_URL = '__SERVER_URL__';

let DATA = [];
let state = [];
let showData = true, showFit = true;

// Format concentration: nM → M scientific notation, 2 decimal places
function formatConc(conc_nM) {
    return (conc_nM * 1e-9).toExponential(2) + ' M';
}
// Round conc to string key for grouping replicates
function concKey(conc_nM) {
    return conc_nM.toPrecision(6);
}
// Build concentration → color mapping so replicates share colors
function buildConcColorMap(cycles) {
    const unique = [];
    cycles.forEach(cy => {
        const k = concKey(cy.conc_nM);
        if (!unique.includes(k)) unique.push(k);
    });
    const map = {};
    unique.forEach((k, i) => { map[k] = COLORS[i % COLORS.length]; });
    return map;
}
function cycleColor(concColorMap, cy) {
    return concColorMap[concKey(cy.conc_nM)] || COLORS[0];
}

// ═══ Upload Logic ═══

const dropZone = document.getElementById('dropZone');
const dataFilesInput = document.getElementById('dataFiles');
const refDropZone = document.getElementById('refDropZone');
const refFileInput = document.getElementById('refFile');

function setupDropZone(zone, input, onUpdate) {
    ['dragenter','dragover'].forEach(e => zone.addEventListener(e, ev => {
        ev.preventDefault(); zone.classList.add('dragover');
    }));
    ['dragleave','drop'].forEach(e => zone.addEventListener(e, ev => {
        ev.preventDefault(); zone.classList.remove('dragover');
    }));
    zone.addEventListener('drop', ev => {
        input.files = ev.dataTransfer.files;
        onUpdate();
    });
    input.addEventListener('change', onUpdate);
}

setupDropZone(dropZone, dataFilesInput, updateFileList);
setupDropZone(refDropZone, refFileInput, updateRefFileList);

function updateFileList() {
    const fl = document.getElementById('fileList');
    const files = dataFilesInput.files;
    if (files.length === 0) { fl.innerHTML = ''; return; }
    fl.innerHTML = Array.from(files).map(f => f.name).join(', ');
    // Auto-set label from first filename
    const labelInput = document.getElementById('dataLabel');
    if (!labelInput.value) {
        labelInput.value = files[0].name.replace(/\.[^.]+$/, '');
    }
}

function updateRefFileList() {
    const fl = document.getElementById('refFileList');
    const files = refFileInput.files;
    if (files.length === 0) { fl.innerHTML = ''; return; }
    fl.innerHTML = 'Reference: ' + files[0].name;
}

async function doUpload() {
    const files = dataFilesInput.files;
    if (files.length === 0) {
        document.getElementById('uploadStatus').textContent = 'Please select at least one data file.';
        return;
    }

    const btn = document.getElementById('uploadBtn');
    const status = document.getElementById('uploadStatus');
    btn.disabled = true;
    btn.textContent = 'Analyzing...';
    status.textContent = 'Uploading and fitting data...';

    const model = document.getElementById('uploadModel').value;
    const label = document.getElementById('dataLabel').value ||
                  files[0].name.replace(/\.[^.]+$/, '');
    const formData = new FormData();
    for (let i = 0; i < files.length; i++) {
        formData.append('data_files', files[i]);
    }
    if (refFileInput.files.length > 0) {
        formData.append('ref_file', refFileInput.files[0]);
    }
    formData.append('model', model);
    formData.append('label', label);

    try {
        const resp = await fetch(SERVER_URL + '/upload', {
            method: 'POST',
            body: formData
        });
        const result = await resp.json();

        if (result.error) {
            status.textContent = 'Error: ' + result.error;
        } else {
            const newResults = result.results || [];
            if (newResults.length === 0) {
                status.textContent = 'No valid results from fitting.';
            } else {
                // Add new results to DATA
                newResults.forEach(r => {
                    DATA.push(r);
                    state.push({
                        excluded: new Set(r.excluded || []),
                        currentData: r
                    });
                });
                status.textContent = newResults.length + ' dataset(s) fitted successfully!';
                showDashboard();
            }
        }
    } catch (e) {
        status.textContent = 'Server error: ' + e.message;
    }

    btn.disabled = false;
    btn.textContent = 'Analyze';
}

function showDashboard() {
    document.getElementById('uploadPanel').style.display = 'none';
    document.getElementById('dashboard').style.display = 'block';
    rebuildSelector();
    buildSummary();
    buildCycleToggles();
    updatePlots();
}

function showUpload() {
    document.getElementById('uploadPanel').style.display = 'block';
    // Don't hide dashboard — allow switching back
    // Clear form
    dataFilesInput.value = '';
    document.getElementById('refFile').value = '';
    document.getElementById('dataLabel').value = '';
    document.getElementById('fileList').innerHTML = '';
    document.getElementById('refFileList').innerHTML = '';
    document.getElementById('uploadStatus').textContent = '';
    window.scrollTo(0, 0);
}

// ═══ Dashboard Logic ═══

function rebuildSelector() {
    const sel = document.getElementById('datasetSelect');
    sel.innerHTML = '';
    DATA.forEach((d, i) => {
        const opt = document.createElement('option');
        opt.value = i;
        opt.textContent = d.analyte || d.filename || ('Dataset ' + (i+1));
        sel.appendChild(opt);
    });
    sel.value = DATA.length - 1;  // Select newest
}

const sumBody = document.getElementById('summaryBody');
function buildSummary() {
    sumBody.innerHTML = '';
    DATA.forEach((d, i) => {
        const s = state[i];
        const cd = s.currentData;
        const tr = document.createElement('tr');
        tr.style.cursor = 'pointer';
        tr.onclick = () => {
            document.getElementById('datasetSelect').value = i;
            onDatasetChange();
        };
        const refLabel = cd.reference ? 'Yes' : 'No';
        const KD = cd.kd / cd.ka;
        const mode = (cd.mode || 'mck').toUpperCase();
        tr.innerHTML = '<td>'+(d.analyte||d.filename||'')+'</td>' +
            '<td>'+mode+'</td>' +
            '<td>'+refLabel+'</td>' +
            '<td class="metric">'+cd.ka.toExponential(3)+'</td>' +
            '<td class="metric">'+cd.kd.toExponential(3)+'</td>' +
            '<td class="metric">'+KD.toExponential(3)+'</td>' +
            '<td class="metric">'+cd.Rmax.toFixed(2)+'</td>' +
            '<td class="metric">'+cd.R2.toFixed(5)+'</td>' +
            '<td>'+(cd.model||'langmuir')+'</td>';
        sumBody.appendChild(tr);
    });
}

function toggleTrace(which) {
    if (which === 'data') showData = !showData;
    if (which === 'fit') showFit = !showFit;
    document.getElementById('btnData').classList.toggle('active', showData);
    document.getElementById('btnFit').classList.toggle('active', showFit);
    updatePlots();
}

function onDatasetChange() {
    buildCycleToggles();
    updatePlots();
}

function buildCycleToggles() {
    const idx = parseInt(document.getElementById('datasetSelect').value);
    if (idx < 0 || idx >= DATA.length) return;
    const d = state[idx].currentData;
    const s = state[idx];
    const container = document.getElementById('cycleToggles');
    container.innerHTML = '';
    const ccm = buildConcColorMap(d.cycles);
    d.cycles.forEach((cy, i) => {
        const color = cycleColor(ccm, cy);
        const label = document.createElement('label');
        label.className = 'cycle-cb' + (s.excluded.has(i) ? ' excluded' : '');
        label.style.borderColor = color;
        const cb = document.createElement('input');
        cb.type = 'checkbox';
        cb.checked = !s.excluded.has(i);
        cb.onchange = () => {
            if (cb.checked) s.excluded.delete(i);
            else s.excluded.add(i);
            label.classList.toggle('excluded', s.excluded.has(i));
        };
        const dot = document.createElement('span');
        dot.style.cssText = 'width:8px;height:8px;border-radius:50%;background:'+color+';display:inline-block;';
        label.appendChild(cb);
        label.appendChild(dot);
        label.appendChild(document.createTextNode(' ' + formatConc(cy.conc_nM)));
        container.appendChild(label);
    });
}

// ── Re-fit via server ──
async function doRefit() {
    const sel = document.getElementById('datasetSelect');
    const idx = parseInt(sel.value);
    const d = DATA[idx];
    const s = state[idx];
    const model = document.getElementById('modelSelect').value;
    const btn = document.getElementById('refitBtn');
    const logEl = document.getElementById('fitLog');

    btn.disabled = true;
    btn.textContent = 'Fitting...';
    logEl.style.display = 'block';
    logEl.textContent = 'Model: ' + model + '\n';
    logEl.textContent += 'Excluded: ' + (s.excluded.size === 0 ? 'none' : [...s.excluded].join(', ')) + '\n';

    try {
        const resp = await fetch(SERVER_URL + '/refit', {
            method: 'POST',
            headers: {'Content-Type': 'application/json'},
            body: JSON.stringify({
                data_file: d.data_file,
                ref_file: d.ref_file_path || null,
                model: model,
                exclude: [...s.excluded],
                ri: document.getElementById('enableRI').checked,
                drift: document.getElementById('enableDrift').checked,
                tc: document.getElementById('enableTC').checked,
            })
        });
        const result = await resp.json();

        if (result.error) {
            logEl.textContent += 'Error: ' + result.error + '\n';
        } else {
            result.analyte = d.analyte;
            result.data_file = d.data_file;
            result.ref_file_path = d.ref_file_path;
            s.currentData = result;

            logEl.textContent += 'ka = ' + result.ka.toExponential(4) + ' M-1s-1\n';
            logEl.textContent += 'kd = ' + result.kd.toExponential(4) + ' s-1\n';
            logEl.textContent += 'KD = ' + (result.kd/result.ka).toExponential(4) + ' M\n';
            logEl.textContent += 'Rmax = ' + result.Rmax.toFixed(3) + ' RU\n';
            if (result.ka2) {
                logEl.textContent += 'ka2 = ' + result.ka2.toExponential(4) + '\n';
                logEl.textContent += 'kd2 = ' + result.kd2.toExponential(4) + '\n';
                if (result.Rmax2) logEl.textContent += 'Rmax2 = ' + result.Rmax2.toFixed(3) + ' RU\n';
            }
            logEl.textContent += 'R\u00b2 = ' + result.R2.toFixed(6) + '  RMS = ' + result.rms.toFixed(4) + '  \u03c7\u00b2 = ' + (result.chi2||0).toExponential(3) + '\n';
            if (result.advanced && (result.advanced.ri || result.advanced.drift || result.advanced.tc))
                logEl.textContent += 'Advanced: ' + (result.advanced.ri ? 'RI ' : '') + (result.advanced.drift ? 'Drift ' : '') + (result.advanced.tc ? 'TC' : '') + ' enabled\n';

            updatePlots();
            buildSummary();
        }
    } catch (e) {
        logEl.textContent += 'Server error: ' + e.message + '\n';
    }

    btn.disabled = false;
    btn.textContent = 'Re-fit';
}

// ── Export PPTX ──
async function doExport() {
    const btn = document.getElementById('exportBtn');
    btn.disabled = true;
    btn.textContent = 'Generating...';

    try {
        const idx = parseInt(document.getElementById('datasetSelect').value);
        const resp = await fetch(SERVER_URL + '/export', {
            method: 'POST',
            headers: {'Content-Type': 'application/json'},
            body: JSON.stringify({ dataset_idx: idx })
        });
        if (!resp.ok) throw new Error('Export failed');
        const blob = await resp.blob();
        const url = URL.createObjectURL(blob);
        const a = document.createElement('a');
        const label = DATA[idx].analyte || 'spr_analysis';
        a.href = url;
        a.download = label + '_SPR.pptx';
        a.click();
        URL.revokeObjectURL(url);
    } catch (e) {
        alert('Export error: ' + e.message);
    }

    btn.disabled = false;
    btn.textContent = 'Export PPTX';
}

function updatePlots() {
    const sel = document.getElementById('datasetSelect');
    const idx = parseInt(sel.value);
    if (idx < 0 || idx >= DATA.length) return;
    const d = state[idx].currentData;
    const s = state[idx];

    // ── Sensorgram ──
    const traces = [];
    const ccm = buildConcColorMap(d.cycles);
    const seenConc = new Set();
    d.cycles.forEach((cy, i) => {
        const skip = cy.skip || [];
        const excluded = s.excluded.has(i);
        const opacity = excluded ? 0.25 : 1;
        const color = cycleColor(ccm, cy);
        const ck = concKey(cy.conc_nM);
        const firstOfConc = !seenConc.has(ck);
        seenConc.add(ck);

        if (showData) {
            traces.push({
                x: cy.time,
                y: cy.response.map((v, j) => (skip[j]) ? null : v),
                mode: 'lines', name: formatConc(cy.conc_nM) + (excluded ? ' (excl)' : ''),
                line: { color: color, width: 1.5 },
                opacity: opacity,
                connectgaps: false,
                legendgroup: ck, showlegend: firstOfConc
            });
        }
        if (showFit) {
            traces.push({
                x: cy.time, y: cy.fitted,
                mode: 'lines', name: formatConc(cy.conc_nM) + ' (fit)',
                line: { color: color, width: 2, dash: 'dash' },
                opacity: excluded ? 0.15 : 1,
                legendgroup: ck, showlegend: false
            });
        }
    });

    const tAssoc = d.t_assoc_end || 60;
    const shapes = [{
        type: 'line', x0: tAssoc, x1: tAssoc, y0: 0, y1: 1,
        yref: 'paper', line: { color: '#bdc3c7', width: 1, dash: 'dot' }
    }];

    const modelLabel = (d.model || 'langmuir') === 'langmuir' ? '1:1 Langmuir' :
                       d.model === 'heterogeneous' ? 'Heterogeneous' : 'Two-State';
    const titleText = (d.analyte || d.filename || '') + ' \u2014 ' + modelLabel;

    Plotly.react('sensorgram', traces, {
        title: { text: titleText, font: { size: 14 } },
        xaxis: { title: 'Time (s)', zeroline: false },
        yaxis: { title: 'Response (RU)', zeroline: true },
        shapes: shapes,
        margin: { l: 60, r: 120, t: 40, b: 50 },
        legend: { x: 1.02, y: 1, xanchor: 'left', yanchor: 'top', font: { size: 9 }, tracegroupgap: 2 },
        hovermode: 'closest'
    }, { responsive: true });

    // ── Residuals ──
    const resTraces = [];
    d.cycles.forEach((cy, i) => {
        const skip = cy.skip || [];
        const excluded = s.excluded.has(i);
        const res = cy.response.map((r, j) => (skip[j] || excluded) ? null : r - cy.fitted[j]);
        resTraces.push({
            x: cy.time, y: res,
            mode: 'lines', name: formatConc(cy.conc_nM),
            line: { color: cycleColor(ccm, cy), width: 1 },
            connectgaps: false, showlegend: false
        });
    });

    Plotly.react('residuals', resTraces, {
        title: { text: 'Residuals (Data \u2212 Fit)', font: { size: 14 } },
        xaxis: { title: 'Time (s)', zeroline: false },
        yaxis: { title: 'Residual (RU)', zeroline: true,
                  zerolinecolor: '#e74c3c', zerolinewidth: 1 },
        shapes: shapes,
        margin: { l: 60, r: 20, t: 40, b: 50 },
        hovermode: 'closest'
    }, { responsive: true });

    // ── Steady-State Affinity ──
    if (d.ss_req && d.concentrations) {
        const concs_nM = [], reqs = [];
        d.concentrations.forEach((c, i) => {
            if (!s.excluded.has(i)) {
                concs_nM.push(c);
                reqs.push(d.ss_req[i]);
            }
        });
        if (concs_nM.length >= 2) {
            const ssKD = d.ss_KD_nM || 1;
            const ssRmax = d.ss_Rmax || 1;
            const minC = Math.min(...concs_nM);
            const maxC = Math.max(...concs_nM) * 2;
            const curveC_M = [], curveR = [];
            const logMin = Math.log10(minC * 1e-9 * 0.3);
            const logMax = Math.log10(maxC * 1e-9);
            for (let i = 0; i <= 200; i++) {
                const c_M = Math.pow(10, logMin + (logMax - logMin) * i / 200);
                curveC_M.push(c_M);
                const c_nM = c_M * 1e9;
                curveR.push(ssRmax * c_nM / (ssKD + c_nM));
            }
            const concs_M = concs_nM.map(c => c * 1e-9);
            Plotly.react('isotherm', [
                { x: concs_M, y: reqs, mode: 'markers', name: 'Req (data)',
                  marker: { size: 10, color: '#e74c3c' } },
                { x: curveC_M, y: curveR, mode: 'lines', name: 'Steady-state fit',
                  line: { color: '#2c3e50', width: 2 } }
            ], {
                title: { text: 'Steady-State Affinity: KD = ' + (ssKD*1e-9).toExponential(3) + ' M, R\u00b2 = ' + (d.ss_R2||0).toFixed(4),
                         font: { size: 13 } },
                xaxis: { title: 'Concentration (M)', type: 'log',
                         exponentformat: 'e' },
                yaxis: { title: 'Req (RU)' },
                margin: { l: 60, r: 20, t: 45, b: 50 },
                showlegend: true,
                legend: { x: 0.98, xanchor: 'right', y: 0.05, yanchor: 'bottom', font: { size: 10 } }
            }, { responsive: true });
        }
    }

    // ── Parameters table ──
    const pb = document.getElementById('paramsBody');
    const KD = d.kd / d.ka;
    let rows = '';
    const modelName = (d.model || 'langmuir');

    if (modelName === 'langmuir') {
        rows = '<tr><td colspan="3" style="background:#ecf0f1;font-weight:600;">1:1 Langmuir</td></tr>' +
            '<tr><td>k<sub>a</sub></td><td class="metric">' + d.ka.toExponential(4) + '</td><td>M<sup>-1</sup>s<sup>-1</sup></td></tr>' +
            '<tr><td>k<sub>d</sub></td><td class="metric">' + d.kd.toExponential(4) + '</td><td>s<sup>-1</sup></td></tr>' +
            '<tr class="highlight"><td>K<sub>D</sub></td><td class="metric">' + KD.toExponential(4) + '</td><td>M</td></tr>' +
            '<tr><td>R<sub>max</sub></td><td class="metric">' + d.Rmax.toFixed(3) + '</td><td>RU</td></tr>';
    } else if (modelName === 'heterogeneous') {
        const KD1 = d.kd / d.ka, KD2 = (d.kd2||1) / (d.ka2||1);
        rows = '<tr><td colspan="3" style="background:#ecf0f1;font-weight:600;">Heterogeneous Ligand</td></tr>' +
            '<tr><td colspan="3" style="background:#f7f9fa;font-style:italic;">Site 1</td></tr>' +
            '<tr><td>k<sub>a1</sub></td><td class="metric">' + d.ka.toExponential(4) + '</td><td>M<sup>-1</sup>s<sup>-1</sup></td></tr>' +
            '<tr><td>k<sub>d1</sub></td><td class="metric">' + d.kd.toExponential(4) + '</td><td>s<sup>-1</sup></td></tr>' +
            '<tr class="highlight"><td>K<sub>D1</sub></td><td class="metric">' + KD1.toExponential(4) + '</td><td>M</td></tr>' +
            '<tr><td>R<sub>max1</sub></td><td class="metric">' + d.Rmax.toFixed(3) + '</td><td>RU</td></tr>' +
            '<tr><td colspan="3" style="background:#f7f9fa;font-style:italic;">Site 2</td></tr>' +
            '<tr><td>k<sub>a2</sub></td><td class="metric">' + (d.ka2||0).toExponential(4) + '</td><td>M<sup>-1</sup>s<sup>-1</sup></td></tr>' +
            '<tr><td>k<sub>d2</sub></td><td class="metric">' + (d.kd2||0).toExponential(4) + '</td><td>s<sup>-1</sup></td></tr>' +
            '<tr class="highlight"><td>K<sub>D2</sub></td><td class="metric">' + KD2.toExponential(4) + '</td><td>M</td></tr>' +
            '<tr><td>R<sub>max2</sub></td><td class="metric">' + (d.Rmax2||0).toFixed(3) + '</td><td>RU</td></tr>';
    } else {
        rows = '<tr><td colspan="3" style="background:#ecf0f1;font-weight:600;">Two-State Conformational Change</td></tr>' +
            '<tr><td colspan="3" style="background:#f7f9fa;font-style:italic;">Step 1: A + B \u21cc AB</td></tr>' +
            '<tr><td>k<sub>a1</sub></td><td class="metric">' + d.ka.toExponential(4) + '</td><td>M<sup>-1</sup>s<sup>-1</sup></td></tr>' +
            '<tr><td>k<sub>d1</sub></td><td class="metric">' + d.kd.toExponential(4) + '</td><td>s<sup>-1</sup></td></tr>' +
            '<tr class="highlight"><td>K<sub>D1</sub></td><td class="metric">' + KD.toExponential(4) + '</td><td>M</td></tr>' +
            '<tr><td colspan="3" style="background:#f7f9fa;font-style:italic;">Step 2: AB \u21cc AB*</td></tr>' +
            '<tr><td>k<sub>a2</sub></td><td class="metric">' + (d.ka2||0).toExponential(4) + '</td><td>s<sup>-1</sup></td></tr>' +
            '<tr><td>k<sub>d2</sub></td><td class="metric">' + (d.kd2||0).toExponential(4) + '</td><td>s<sup>-1</sup></td></tr>' +
            '<tr><td>R<sub>max</sub></td><td class="metric">' + d.Rmax.toFixed(3) + '</td><td>RU</td></tr>';
    }

    /* ── Quality Control Section ── */
    rows += '<tr><td colspan="3" style="background:#2c3e50;color:#fff;font-weight:600;letter-spacing:0.5px;">Quality Control</td></tr>';

    /* Helper for badge */
    function qcBadge(label, color) {
        return ' <span style="display:inline-block;padding:1px 6px;border-radius:3px;font-size:10px;font-weight:600;color:#fff;background:' + color + ';">' + label + '</span>';
    }

    /* Overall fit quality (R²) */
    {
        let fLabel, fColor;
        if (d.R2 >= 0.99) { fLabel = 'Excellent'; fColor = '#27ae60'; }
        else if (d.R2 >= 0.95) { fLabel = 'Good'; fColor = '#2980b9'; }
        else if (d.R2 >= 0.90) { fLabel = 'Fair'; fColor = '#f39c12'; }
        else { fLabel = 'Poor'; fColor = '#e74c3c'; }
        rows += '<tr><td>R\u00b2</td><td class="metric">' + d.R2.toFixed(6) + qcBadge(fLabel, fColor) + '</td><td></td></tr>';
    }
    rows += '<tr><td>\u03c7\u00b2</td><td class="metric">' + (d.chi2||0).toExponential(4) + '</td><td>RU\u00b2</td></tr>' +
            '<tr><td>RMS</td><td class="metric">' + d.rms.toFixed(4) + '</td><td>RU</td></tr>' +
            '<tr><td>Data points</td><td class="metric">' + d.n_points + '</td><td></td></tr>';

    /* U-value (parameter uniqueness) */
    if (d.u_value !== undefined) {
        const uv = d.u_value;
        let uLabel, uColor;
        if (uv < 15) { uLabel = 'Unique'; uColor = '#27ae60'; }
        else if (uv < 25) { uLabel = 'Borderline'; uColor = '#f39c12'; }
        else { uLabel = 'Correlated'; uColor = '#e74c3c'; }
        rows += '<tr><td>U-value</td><td class="metric">' + uv.toFixed(1) + qcBadge(uLabel, uColor) + '</td><td></td></tr>';
    }

    /* Mass transport assessment */
    if (d.advanced && d.advanced.tc && d.tc !== undefined && d.tc > 0) {
        const ka_lin = d.ka;
        const mtIndex = ka_lin * d.Rmax / d.tc;
        let mtLabel, mtColor;
        if (mtIndex < 0.2) { mtLabel = 'Negligible'; mtColor = '#27ae60'; }
        else if (mtIndex < 1.0) { mtLabel = 'Moderate'; mtColor = '#f39c12'; }
        else { mtLabel = 'Significant'; mtColor = '#e74c3c'; }
        rows += '<tr><td>Mass Transport</td><td class="metric">k<sub>a</sub>\u00b7R<sub>max</sub>/t<sub>c</sub> = ' +
                mtIndex.toFixed(2) + qcBadge(mtLabel, mtColor) + '</td><td></td></tr>' +
                '<tr><td style="padding-left:16px;font-size:11px;color:#777;">t<sub>c</sub></td><td class="metric" style="font-size:11px;">' +
                d.tc.toExponential(3) + '</td><td style="font-size:11px;">RU\u00b7M\u207b\u00b9\u00b7s\u207b\u00b9</td></tr>';
    }

    /* Drift assessment */
    if (d.advanced && d.advanced.drift && d.drift !== undefined) {
        const driftRate = Math.abs(d.drift);
        /* Compare drift contribution over assoc time to Rmax */
        const driftContrib = driftRate * (d.t_assoc_end || 60);
        const driftPct = d.Rmax > 0 ? (driftContrib / d.Rmax * 100) : 0;
        let drLabel, drColor;
        if (driftPct < 2) { drLabel = 'Negligible'; drColor = '#27ae60'; }
        else if (driftPct < 10) { drLabel = 'Moderate'; drColor = '#f39c12'; }
        else { drLabel = 'Significant'; drColor = '#e74c3c'; }
        rows += '<tr><td>Drift</td><td class="metric">' + d.drift.toExponential(3) +
                ' RU/s (' + driftPct.toFixed(1) + '% of R<sub>max</sub>)' +
                qcBadge(drLabel, drColor) + '</td><td></td></tr>';
    }

    /* Bulk RI assessment */
    if (d.advanced && d.advanced.ri) {
        let maxRI = 0;
        d.cycles.forEach((cy) => { if (!cy.excluded && Math.abs(cy.ri) > maxRI) maxRI = Math.abs(cy.ri); });
        const riPct = d.Rmax > 0 ? (maxRI / d.Rmax * 100) : 0;
        let riLabel, riColor;
        if (riPct < 5) { riLabel = 'Negligible'; riColor = '#27ae60'; }
        else if (riPct < 20) { riLabel = 'Moderate'; riColor = '#f39c12'; }
        else { riLabel = 'Significant'; riColor = '#e74c3c'; }
        rows += '<tr><td>Bulk RI</td><td class="metric">max |RI| = ' + maxRI.toFixed(3) +
                ' RU (' + riPct.toFixed(1) + '% of R<sub>max</sub>)' +
                qcBadge(riLabel, riColor) + '</td><td></td></tr>';
        /* Per-cycle RI detail (collapsible sub-rows) */
        d.cycles.forEach((cy, i) => {
            if (cy.excluded) return;
            rows += '<tr><td style="padding-left:16px;font-size:11px;color:#777;">' + formatConc(cy.conc_nM) + '</td>' +
                    '<td class="metric" style="font-size:11px;">RI = ' + cy.ri.toFixed(3) + '</td><td style="font-size:11px;">RU</td></tr>';
        });
    }

    if (d.ss_KD_nM > 0) {
        rows += '<tr><td colspan="3" style="background:#ecf0f1;font-weight:600;">Steady-State Affinity</td></tr>' +
            '<tr class="highlight"><td>K<sub>D</sub> (SS)</td><td class="metric">' + (d.ss_KD_nM*1e-9).toExponential(4) + '</td><td>M</td></tr>' +
            '<tr><td>R<sub>max</sub> (SS)</td><td class="metric">' + d.ss_Rmax.toFixed(3) + '</td><td>RU</td></tr>' +
            '<tr><td>R\u00b2 (SS)</td><td class="metric">' + (d.ss_R2||0).toFixed(6) + '</td><td></td></tr>';
    }
    pb.innerHTML = rows;
}
</script>
</body>
</html>"""


def build_html(server_url=""):
    """Generate HTML dashboard string."""
    html = HTML_TEMPLATE.replace('__COLORS_JSON__', json.dumps(CYCLE_COLORS))
    html = html.replace('__SERVER_URL__', server_url)
    return html


# ── HTTP Server ────────────────────────────────────────────────────

class SPRHandler(BaseHTTPRequestHandler):
    html_content = ""
    datasets = []  # Shared state: all fitted datasets

    def log_message(self, format, *args):
        pass

    def do_GET(self):
        path = urlparse(self.path).path
        if path == "/" or path == "/index.html":
            self.send_response(200)
            self.send_header("Content-Type", "text/html; charset=utf-8")
            self.end_headers()
            self.wfile.write(SPRHandler.html_content.encode())
        else:
            self.send_error(404)

    def do_POST(self):
        path = urlparse(self.path).path
        content_len = int(self.headers.get("Content-Length", 0))
        body = self.rfile.read(content_len)

        if path == "/upload":
            result = self.handle_upload(body)
            self._json_response(result)
        elif path == "/refit":
            result = self.handle_refit(json.loads(body))
            self._json_response(result)
        elif path == "/export":
            self.handle_export(json.loads(body))
        else:
            self.send_error(404)

    def do_OPTIONS(self):
        self.send_response(200)
        self.send_header("Access-Control-Allow-Origin", "*")
        self.send_header("Access-Control-Allow-Methods", "POST, OPTIONS")
        self.send_header("Access-Control-Allow-Headers", "Content-Type")
        self.end_headers()

    def _json_response(self, data):
        self.send_response(200)
        self.send_header("Content-Type", "application/json")
        self.send_header("Access-Control-Allow-Origin", "*")
        self.end_headers()
        self.wfile.write(json.dumps(data).encode())

    def handle_upload(self, body):
        content_type = self.headers.get("Content-Type", "")
        parts = parse_multipart(content_type, body)

        if 'data_files' not in parts:
            return {"error": "No data files uploaded"}

        model = parts.get('model', 'langmuir')
        label = parts.get('label', 'Dataset')

        # Handle single or multiple files
        data_files = parts['data_files']
        if isinstance(data_files, dict):
            data_files = [data_files]
        elif isinstance(data_files, list):
            pass
        else:
            data_files = [data_files]

        # Save uploaded files to temp directory
        ref_path = None
        if 'ref_file' in parts and isinstance(parts['ref_file'], dict):
            ref_info = parts['ref_file']
            ref_path = os.path.join(UPLOAD_DIR, 'ref_' + ref_info['filename'])
            with open(ref_path, 'wb') as f:
                f.write(ref_info['data'])
            print(f"  Reference file: {ref_info['filename']}")

        results = []
        for i, file_info in enumerate(data_files):
            if not isinstance(file_info, dict):
                continue
            fname = file_info['filename']
            fpath = os.path.join(UPLOAD_DIR, fname)
            with open(fpath, 'wb') as f:
                f.write(file_info['data'])

            file_label = label if len(data_files) == 1 else f"{label} ({fname})"
            print(f"\nProcessing: {fname}")
            result = run_fitter(fpath, ref_path, model)
            if result:
                result['analyte'] = file_label
                result['data_file'] = fpath
                result['ref_file_path'] = ref_path
                SPRHandler.datasets.append(result)
                results.append(result)
            else:
                print(f"  Failed to fit {fname}")

        return {"results": results}

    def handle_refit(self, body):
        data_file = body.get("data_file")
        ref_file = body.get("ref_file")
        model = body.get("model", "langmuir")
        exclude = body.get("exclude", [])
        ri = body.get("ri", False)
        drift = body.get("drift", False)
        tc = body.get("tc", False)

        if not data_file or not os.path.exists(data_file):
            return {"error": f"Data file not found: {data_file}"}

        print(f"  Re-fit: {os.path.basename(data_file)} model={model} exclude={exclude} ri={ri} drift={drift} tc={tc}")
        result = run_fitter(data_file, ref_file, model, exclude, ri=ri, drift=drift, tc=tc)
        if result is None:
            return {"error": "Fitting failed"}
        return result

    def handle_export(self, body):
        idx = body.get("dataset_idx", None)
        print(f"  Export PPTX: dataset {idx}")

        try:
            pptx_bytes = generate_pptx(SPRHandler.datasets, idx)
            self.send_response(200)
            self.send_header("Content-Type",
                           "application/vnd.openxmlformats-officedocument.presentationml.presentation")
            self.send_header("Content-Length", str(len(pptx_bytes)))
            self.send_header("Access-Control-Allow-Origin", "*")
            self.end_headers()
            self.wfile.write(pptx_bytes)
        except Exception as e:
            import traceback
            traceback.print_exc()
            self._json_response({"error": str(e)})


def main():
    import argparse
    parser = argparse.ArgumentParser(description="SPR Analysis Server")
    parser.add_argument("--port", type=int, default=8765)
    args = parser.parse_args()

    # Create temp upload directory
    global UPLOAD_DIR
    UPLOAD_DIR = tempfile.mkdtemp(prefix="spr_uploads_")
    print(f"Upload directory: {UPLOAD_DIR}")

    # Compile
    if not compile_fitter():
        sys.exit(1)

    # Build HTML
    server_url = f"http://localhost:{args.port}"
    SPRHandler.html_content = build_html(server_url)

    # Start server
    server = HTTPServer(("", args.port), SPRHandler)
    print(f"\nServer running at {server_url}")
    print("Open in browser to upload and analyze SPR data.")
    print("Press Ctrl+C to stop.\n")

    try:
        server.serve_forever()
    except KeyboardInterrupt:
        print("\nServer stopped.")
        server.server_close()
        # Clean up temp files
        import shutil
        shutil.rmtree(UPLOAD_DIR, ignore_errors=True)


if __name__ == "__main__":
    main()
